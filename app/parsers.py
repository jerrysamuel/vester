import os
from PyPDF2 import PdfReader
from pptx import Presentation
from pdf2image import convert_from_path
import pytesseract

def parse_pdf(filepath):
    """Parse text from PDF using PyPDF2 or fallback to pytesseract."""
    content = []
    try:
        reader = PdfReader(filepath)
        for page in reader.pages:
            text = page.extract_text()
            if text:
                content.append(text.strip())
        if not content:
            content = ocr_pdf(filepath)
    except Exception:
        content = ocr_pdf(filepath)
    return " ".join(content)

def ocr_pdf(filepath):
    """OCR for PDFs."""
    images = convert_from_path(filepath)
    return " ".join([pytesseract.image_to_string(image).strip() for image in images])

def parse_pptx(filepath):
    """Parse text from PPTX using python-pptx."""
    content = []
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content.append(shape.text.strip())
    except Exception:
        pass
    return " ".join(content)

def parse_file(filepath):
    """Determine file type and parse."""
    _, ext = os.path.splitext(filepath)
    if ext.lower() == '.pdf':
        return parse_pdf(filepath)
    elif ext.lower() == '.pptx':
        return parse_pptx(filepath)
    else:
        raise ValueError("Unsupported file format.")
